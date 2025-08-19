# app.py

import io

import dash
import dash_daq as daq
import numpy as np
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
from dash import ALL, Input, Output, State, dcc, html
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
# PDF export (with embedded chart images)
from reportlab.platypus import Image as RLImage
from reportlab.platypus import (ListFlowable, ListItem, Paragraph,
                                SimpleDocTemplate, Spacer, Table, TableStyle)

from config import DOMAINS, FRAMEWORK_ROWS, LIKERT, QUESTIONS, RECS

app = dash.Dash(__name__, suppress_callback_exceptions=True)
app.title = "Maturity & Governance Assessment"
server = app.server


def _validate_config() -> None:
    """
    Check the sanity of the configuration in `config.py`.

    Prints warnings if questions have unknown domains or non-list tags.
    """
    bad_domains = [
        q for q in QUESTIONS if q["domain"] not in DOMAINS
    ]  # unknown domains
    bad_tags = [
        q for q in QUESTIONS if not isinstance(q.get("tags", []), (list, tuple))
    ]  # non-list tags

    if bad_domains or bad_tags:
        print("[config warning]")
        if bad_domains:
            print(
                "  Questions with unknown domain:",
                [q["id"] for q in bad_domains],
            )
        if bad_tags:
            print(
                "  Questions with non-list tags:",
                [q["id"] for q in bad_tags],
            )


_validate_config()


# ----------- Helpers -------------
def likert_to_pct(v):
    """
    Convert a Likert scale value (1-5) to a percentage.

    1  -> 0%
    2  -> 25%
    3  -> 50%
    4  -> 75%
    5  -> 100%
    """

    return round((int(v) - 1) / 4 * 100, 2)


def pct_to_level(p):
    """
    Convert a percentage to a level (1-5).

    0%  -> 1
    25% -> 2
    50% -> 3
    75% -> 4
    100%-> 5

    :param p: percentage (0-100)
    :return: level (1-5)
    """
    return int(min(5, max(1, round((p / 100) * 4 + 1))))


def _as_list(x):
    """
    Convert the input to a list.

    If the input is already a list or tuple, return it as is.
    If the input is None or pd.NA, return an empty list.
    Otherwise, wrap the input in a list.

    :param x: input to convert to a list
    :return: a list
    """
    return (
        list(x) if isinstance(x, (list, tuple)) else ([] if x in (None, pd.NA) else [x])
    )


def _slug(s: str):
    return "".join(ch.lower() if ch.isalnum() else "-" for ch in s).strip("-")


# for chart sizes
RADAR_H = 360
BAR_H = 360
HEAT_H = 520


def _base_fig_layout(fig, theme="light", height=360):
    """
    Apply a consistent layout to a figure.

    This sets the font to a contrasting color for light/dark themes,
    and sets the grid color to a contrasting color. It also sets the
    axis colors to match the text color.

    :param fig: a figure to update
    :param theme: a string, either "light" or "dark"
    :param height: the height of the figure in pixels
    :return: the updated figure
    """

    font_color = "#f6f7fb" if theme == "dark" else "#0b1020"
    grid_color = "#334155" if theme == "dark" else "#CBD5E1"
    axis_color = font_color
    fig.update_layout(
        autosize=False,
        height=height,
        margin=dict(l=30, r=30, t=30, b=30),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(color=font_color),
        xaxis=dict(
            showgrid=True,
            gridcolor=grid_color,
            zeroline=False,
            linecolor=axis_color,
            ticks="outside",
            fixedrange=True,
        ),
        yaxis=dict(
            showgrid=True,
            gridcolor=grid_color,
            zeroline=False,
            linecolor=axis_color,
            ticks="outside",
            fixedrange=True,
        ),
        uirevision="keep",
    )
    return fig


# -------------- Layout --------------------
def build_question_cards():
    """
    Build a list of HTML cards, one per domain, each containing a sequence of
    questions with Likert scale input controls.

    :return: a list of HTML Div elements, each representing a domain card
    """
    groups = {}
    for q in QUESTIONS:
        groups.setdefault(q["domain"], []).append(q)
    cards = []
    for domain in DOMAINS:
        qlist = groups.get(domain, [])
        children = [html.H3(domain, className="domain-title")]
        for q in qlist:
            rid = {"type": "q-input", "qid": q["id"]}
            children.append(
                html.Div(
                    [
                        html.Div(q["text"], className="qtext"),
                        # vertical Likert handled by CSS
                        dcc.RadioItems(
                            id=rid, options=LIKERT, value=3, className="likert"
                        ),
                    ],
                    className="qrow",
                )
            )
        cards.append(html.Div(children, className=f"domain-card d-{_slug(domain)}"))
    return cards


app.layout = html.Div(
    id="page-root",
    className="page theme-light",
    children=[
        dcc.Store(id="responses-store"),
        dcc.Store(id="questions-store", data=QUESTIONS),
        dcc.Store(id="theme-store", data="light"),
        # Header
        html.Div(
            [
                html.H1("Data Maturity & Governance Assessment"),
                html.Div(
                    [
                        html.Div(
                            [
                                html.Label("Organization"),
                                dcc.Input(
                                    id="org-name",
                                    placeholder="e.g., AllureAfrica",
                                    className="textin",
                                ),
                            ],
                            className="field",
                        ),
                        html.Div(
                            [
                                html.Label("Assessor"),
                                dcc.Input(
                                    id="assessor",
                                    placeholder="Your name",
                                    className="textin",
                                ),
                            ],
                            className="field",
                        ),
                        html.Div(
                            [
                                html.Label("Sector"),
                                dcc.Input(
                                    id="sector",
                                    placeholder="e.g., Fintech",
                                    className="textin",
                                ),
                            ],
                            className="field",
                        ),
                        html.Div(
                            [
                                html.Label("Dark mode"),
                                daq.BooleanSwitch(
                                    id="theme-switch",
                                    on=False,
                                    color="#4f46e5",
                                    className="theme-switch",
                                ),
                            ],
                            className="field",
                        ),
                    ],
                    className="meta",
                ),
            ],
            className="header",
        ),
        dcc.Tabs(
            id="tabs",
            value="tab-assess",
            children=[
                dcc.Tab(
                    label="Assessment",
                    value="tab-assess",
                    children=[
                        html.Div(build_question_cards(), className="grid"),
                        html.Button(
                            "Compute Scores",
                            id="submit-assessment",
                            n_clicks=0,
                            className="primary",
                        ),
                    ],
                ),
                dcc.Tab(
                    label="Results & Insights",
                    value="tab-results",
                    children=[
                        html.Div(id="kpis", className="kpis"),
                        # Export controls
                        html.Div(
                            [
                                html.Button(
                                    "Download CSV",
                                    id="dl-csv",
                                    n_clicks=0,
                                    className="secondary",
                                ),
                                dcc.Download(id="dl-csv-out"),
                                html.Button(
                                    "Download PPTX",
                                    id="dl-ppt",
                                    n_clicks=0,
                                    className="secondary",
                                ),
                                dcc.Download(id="dl-ppt-out"),
                                html.Button(
                                    "Download PDF",
                                    id="dl-pdf",
                                    n_clicks=0,
                                    className="secondary",
                                ),
                                dcc.Download(id="dl-pdf-out"),
                            ],
                            className="export-row",
                        ),
                        # Charts row (fixed heights)
                        html.Div(
                            [
                                dcc.Graph(
                                    id="radar",
                                    style={"height": f"{RADAR_H}px"},
                                    config={
                                        "responsive": False,
                                        "displaylogo": False,
                                        "scrollZoom": False,
                                    },
                                ),
                                dcc.Graph(
                                    id="bar",
                                    style={"height": f"{BAR_H}px"},
                                    config={
                                        "responsive": False,
                                        "displaylogo": False,
                                        "scrollZoom": False,
                                    },
                                ),
                            ],
                            className="charts",
                        ),
                        # Heatmap + recommendations
                        html.Div(
                            className="row-heat-actions",
                            children=[
                                html.Div(
                                    [
                                        html.H3("Control Coverage Heatmap"),
                                        dcc.Graph(
                                            id="heatmap",
                                            style={"height": f"{HEAT_H}px"},
                                            config={
                                                "responsive": False,
                                                "displaylogo": False,
                                                "scrollZoom": False,
                                            },
                                        ),
                                    ],
                                    className="col heatmap-col",
                                ),
                                html.Div(
                                    [
                                        html.H3("Top Recommended Actions"),
                                        html.Ul(id="actions-list", className="actions"),
                                    ],
                                    className="col recs-col",
                                ),
                            ],
                        ),
                    ],
                ),
            ],
        ),
    ],
)


# -------------- Scoring & Aggregation ---------------
def compute_scores(resp_df):
    """
    Compute overall and domain-level scores from a response dataframe.

    Args:
        resp_df (pd.DataFrame): response dataframe with columns "value" and "weight"

    Returns:
        tuple: (domain_scores, overall)
            domain_scores (dict): mapping of domain to score
            overall (float): mean of domain scores, rounded to 2 decimal places
    """
    resp_df = resp_df.copy()
    resp_df["pct"] = resp_df["value"].apply(likert_to_pct)
    resp_df["wx"] = resp_df["pct"] * resp_df["weight"]
    dgroups = resp_df.groupby("domain", as_index=False).agg(
        score=("wx", "sum"), w=("weight", "sum")
    )
    dgroups["score"] = (dgroups["score"] / dgroups["w"]).round(2)
    domain_scores = dict(zip(dgroups["domain"], dgroups["score"]))
    overall = round(dgroups["score"].mean(), 2) if not dgroups.empty else 0.0
    return domain_scores, overall


def control_coverage_matrix(resp_df):
    """
    Generate a dataframe with one row per (Framework, Domain) combination, with a "Coverage" column indicating the mean control coverage score for that combination.

    Args:
        resp_df (pd.DataFrame): response dataframe with columns "domain", "value", and "tags"

    Returns:
        pd.DataFrame: coverage matrix
    """
    resp_df = resp_df.copy()
    resp_df["tags"] = resp_df["tags"].apply(_as_list)
    rows = []
    for fw in FRAMEWORK_ROWS:
        for dom in DOMAINS:
            g = resp_df[
                (resp_df["domain"] == dom) & (resp_df["tags"].apply(lambda t: fw in t))
            ]
            val = None if g.empty else round(g["value"].apply(likert_to_pct).mean(), 2)
            rows.append({"Framework": fw, "Domain": dom, "Coverage": val})
    return pd.DataFrame(rows)


def recommendations(domain_scores):
    """
    Return a list of at most 10 action items, sorted by score (ascending),
    where each item is a dict with keys "domain", "score", and "action".
    The items are generated by looking up the top 2 recommended actions for
    each domain in the input domain_scores, and ranking them by score.
    If domain_scores is empty, returns an empty list.
    """
    if not domain_scores:
        return []
    items = []
    for dom, score in domain_scores.items():
        level = pct_to_level(score)
        for a in RECS.get(dom, {}).get(level, [])[:2]:
            items.append({"domain": dom, "score": score, "action": a})
    return sorted(items, key=lambda x: x["score"])[:10]


# ---------- Figures (fixed sizes, consistent) ------------------
def radar_figure(domain_scores, theme="light"):
    """
    Return a radar figure with a fixed size and consistent layout.

    Args:
        domain_scores (dict): mapping of domain to score
        theme (str, optional): light or dark. Defaults to "light".

    Returns:
        go.Figure: radar figure
    """
    vals = [float(domain_scores.get(d, 0.0)) for d in DOMAINS]
    cats = DOMAINS
    cats2, vals2 = cats + [cats[0]], vals + [vals[0]]

    grid_color = "#334155" if theme == "dark" else "#CBD5E1"
    fig = go.Figure()
    fig.add_trace(
        go.Scatterpolar(
            r=vals2,
            theta=cats2,
            fill="toself",
            name="Maturity",
            line=dict(width=2),
            marker=dict(size=4),
            cliponaxis=True,
        )
    )
    fig.update_layout(
        autosize=False,
        height=RADAR_H,  # keep it fixed
        polar=dict(
            radialaxis=dict(
                range=[0, 100],  # pin the range
                autorange=False,  # and stop autoranging
                tick0=0,
                dtick=20,
                gridcolor=grid_color,
                showline=True,
                linewidth=1,
            ),
            angularaxis=dict(gridcolor=grid_color, showline=True, linewidth=1),
        ),
        uirevision="keep",  # stop Plotly from “helpfully” resizing
    )
    return _base_fig_layout(fig, theme, height=RADAR_H)


def bar_figure(domain_scores, theme="light"):
    """
    Return a bar figure with a fixed size and consistent layout.

    Args:
        domain_scores (dict): mapping of domain to score
        theme (str, optional): light or dark. Defaults to "light".

    Returns:
        go.Figure: bar figure
    """
    vals = [float(domain_scores.get(d, 0.0)) for d in DOMAINS]
    fig = go.Figure(go.Bar(x=DOMAINS, y=vals))
    fig.update_layout(
        autosize=False,
        height=BAR_H,
        xaxis=dict(categoryorder="array", categoryarray=DOMAINS, fixedrange=True),
        yaxis=dict(range=[0, 100], fixedrange=True, tick0=0, dtick=20),
        uirevision="keep",
    )
    return _base_fig_layout(fig, theme, height=BAR_H)


def heatmap_figure(mat_df, theme="light"):
    """
    Return a heatmap figure with a fixed size and consistent layout.

    Args:
        mat_df (pd.DataFrame): coverage matrix dataframe
        theme (str, optional): light or dark. Defaults to "light".

    Returns:
        go.Figure: heatmap figure
    """
    if mat_df is None or mat_df.empty:
        pv = pd.DataFrame(index=FRAMEWORK_ROWS, columns=DOMAINS, dtype=float)
    else:
        pv = mat_df.pivot_table(
            index="Framework", columns="Domain", values="Coverage", aggfunc="mean"
        ).reindex(index=FRAMEWORK_ROWS, columns=DOMAINS)
    pv = pv.astype(float)
    z = pv.to_numpy()

    font_color = "#f6f7fb" if theme == "dark" else "#0b1020"
    muted = "#a9b0c4" if theme == "dark" else "#60646e"

    all_nan = np.all(np.isnan(z))
    if all_nan:
        z_display = np.zeros_like(z, dtype=float)
        showscale = False
        annotations = [
            dict(
                text="No tagged coverage yet",
                xref="paper",
                yref="paper",
                x=0.5,
                y=0.5,
                showarrow=False,
                font=dict(size=14, color=muted),
            )
        ]
        colorscale = (
            [[0, "#d8dde9"], [1, "#d8dde9"]]
            if theme == "light"
            else [[0, "#2a334f"], [1, "#2a334f"]]
        )
    else:
        z_display = np.where(np.isnan(z), np.nan, z)
        showscale = True
        annotations = []
        for i, fw in enumerate(pv.index):
            for j, dom in enumerate(pv.columns):
                val = pv.iloc[i, j]
                if pd.notna(val):
                    annotations.append(
                        dict(
                            x=dom,
                            y=fw,
                            text=f"{val:.0f}%",
                            showarrow=False,
                            font=dict(size=11, color=font_color),
                        )
                    )
        colorscale = "Viridis"

    fig = go.Figure(
        data=go.Heatmap(
            z=z_display,
            x=list(pv.columns),
            y=list(pv.index),
            zmin=0,
            zmax=100,
            colorscale=colorscale,
            showscale=showscale,
            hovertemplate="Domain: %{x}<br>Framework: %{y}<br>Coverage: %{z:.0f}%<extra></extra>",
            xgap=1,
            ygap=1,
        )
    )
    fig.update_layout(
        autosize=False,
        height=HEAT_H,
        xaxis=dict(title="", tickangle=0),
        yaxis=dict(title=""),
    )
    return _base_fig_layout(fig, theme, height=HEAT_H)


# -------- Callbacks ------------------
@app.callback(
    Output("responses-store", "data"),
    Input("submit-assessment", "n_clicks"),
    State("questions-store", "data"),
    State("org-name", "value"),
    State("assessor", "value"),
    State("sector", "value"),
    State({"type": "q-input", "qid": ALL}, "value"),
    prevent_initial_call=True,
)
def on_submit(_, qbank, org, assessor, sector, values):
    """
    Stores the user's responses in the "responses-store" store.

    Arguments:
        n_clicks (int): Click count of the submit button.
        qbank (list): The question bank, as stored in the "questions-store".
        org (str): The organization name, as entered in the "org-name" input field.
        assessor (str): The assessor name, as entered in the "assessor" input field.
        sector (str): The sector name, as entered in the "sector" input field.
        values (list): The user's responses, as entered in the question input fields.

    Returns:
        dict: A dict containing the user's responses, domain scores, overall maturity score,
            control coverage matrix, and recommendations.
    """
    qbank = qbank or []
    vals = values or []
    rows = []
    for i, q in enumerate(qbank):
        v = 3
        if (
            i < len(vals)
            and isinstance(vals[i], (int, float))
            and 1 <= int(vals[i]) <= 5
        ):
            v = int(vals[i])
        rows.append(
            {
                "id": q["id"],
                "domain": q["domain"],
                "text": q["text"],
                "weight": float(q["weight"]),
                "tags": q.get("tags", []),
                "value": v,
            }
        )
    df = pd.DataFrame(rows)
    d_scores, overall = compute_scores(df)
    mat = control_coverage_matrix(df)
    recs = recommendations(d_scores)

    return {
        "org": org or "",
        "assessor": assessor or "",
        "sector": sector or "",
        "responses": df.to_dict(orient="records"),
        "domain_scores": d_scores,
        "overall": overall,
        "coverage_matrix": mat.to_dict(orient="records"),
        "recommendations": recs,
    }


@app.callback(
    Output("kpis", "children"),
    Output("radar", "figure"),
    Output("bar", "figure"),
    Output("heatmap", "figure"),
    Output("actions-list", "children"),
    Input("responses-store", "data"),
    Input("theme-store", "data"),
    prevent_initial_call=True,
)
def update_results(data, theme):
    """
    Updates the KPIs, radar, bar, heatmap, and recommendations based on the user's responses.

    Args:
        data (dict): The user's responses, as stored in the "responses-store".
        theme (str): The theme name ("light" or "dark"), as stored in the "theme-store".

    Returns:
        tuple: A tuple containing the updated KPIs, radar, bar, heatmap, and recommendations.
    """
    if not data:
        raise dash.exceptions.PreventUpdate

    d_scores = data.get("domain_scores", {}) or {}
    overall = data.get("overall", 0.0) or 0.0
    mat = pd.DataFrame(data.get("coverage_matrix", []) or [])
    recs = data.get("recommendations", []) or []

    kpi_children = [
        html.Div(
            [
                html.Div("Overall Maturity", className="kpi-title"),
                html.Div(f"{float(overall):.1f}%", className="kpi-value"),
            ],
            className="kpi",
        ),
    ]
    for dom in DOMAINS:
        val = float(d_scores.get(dom, 0.0))
        kpi_children.append(
            html.Div(
                [
                    html.Div(dom, className="kpi-title"),
                    html.Div(f"{val:.1f}%", className="kpi-value"),
                ],
                className="kpi",
            )
        )

    return (
        kpi_children,
        radar_figure(d_scores, theme),
        bar_figure(d_scores, theme),
        heatmap_figure(mat, theme),
        [
            html.Li(f"[{r['domain']}] {r['action']} (current: {r['score']:.0f}%)")
            for r in recs
        ],
    )


# Exports
@app.callback(
    Output("dl-csv-out", "data"),
    Input("dl-csv", "n_clicks"),
    State("responses-store", "data"),
    prevent_initial_call=True,
)
def download_csv(_, data):
    """
    Download the user's responses as a CSV file.

    Args:
        _ (int): Click count of the "Download CSV" button.
        data (dict): The user's responses, as stored in the "responses-store".

    Returns:
        dcc.SendData: A dcc.SendData object containing the CSV data.
    """
    if not data:
        raise dash.exceptions.PreventUpdate
    df = pd.DataFrame(data["responses"])
    df = df[["id", "domain", "text", "weight", "value"]]
    return dcc.send_data_frame(df.to_csv, "data_maturity_responses.csv", index=False)


def _write_ppt_bytes(buf, data):
    """
    Write a PowerPoint presentation with the following slides to a bytes buffer.

    1. Title slide with organization, assessor, and sector information.
    2. Summary slide with overall maturity score, method, and domain information.
    3. Domain scores table.
    4. Top recommended actions table.
    5. Control coverage matrix (framework × domain).

    Args:
        buf (BytesIO): A BytesIO object to write the presentation to.
        data (dict): The user's responses and computed results.

    Returns:
        None
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Data Maturity & Governance Assessment"
    slide.placeholders[1].text = (
        f"Organization: {data.get('org','')}\n"
        f"Assessor: {data.get('assessor','')}\n"
        f"Sector: {data.get('sector','')}"
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = f"Overall Maturity: {data.get('overall', 0):.1f}%"
    body.add_paragraph().text = "Method: 18 questions (Likert 1–5) weighted per domain."
    body.add_paragraph().text = (
        "Domains: Governance, Quality, Metadata, Privacy/Sec, Architecture, AI Gov."
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Domain Scores"
    d_scores = data.get("domain_scores", {})
    rows, cols = len(DOMAINS) + 1, 2
    table = slide.shapes.add_table(
        rows, cols, Inches(0.8), Inches(1.5), Inches(8.0), Inches(0.8 + 0.35 * rows)
    ).table
    table.cell(0, 0).text, table.cell(0, 1).text = "Domain", "Score (%)"
    for i, dom in enumerate(DOMAINS, start=1):
        table.cell(i, 0).text = dom
        table.cell(i, 1).text = f"{float(d_scores.get(dom, 0.0)):.1f}"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Recommended Actions"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for rec in data.get("recommendations", []):
        tf.add_paragraph().text = (
            f"[{rec['domain']}] {rec['action']} (current: {rec['score']:.0f}%)"
        )

    # Coverage table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Control Coverage (Framework × Domain)"
    mat = pd.DataFrame(data.get("coverage_matrix", []))
    if mat.empty:
        mat = pd.DataFrame(
            [
                {"Framework": fw, "Domain": dom, "Coverage": None}
                for fw in FRAMEWORK_ROWS
                for dom in DOMAINS
            ]
        )
    pv = mat.pivot_table(
        index="Framework", columns="Domain", values="Coverage", aggfunc="mean"
    ).reindex(index=FRAMEWORK_ROWS, columns=DOMAINS)
    rows, cols = len(pv.index) + 1, len(pv.columns) + 1
    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.4), Inches(9.0), Inches(0.6 + 0.3 * rows)
    ).table
    table.cell(0, 0).text = "Framework"
    for j, dom in enumerate(pv.columns, start=1):
        table.cell(0, j).text = dom
    for i, fw in enumerate(pv.index, start=1):
        table.cell(i, 0).text = fw
        for j, dom in enumerate(pv.columns, start=1):
            val = pv.loc[fw, dom]
            table.cell(i, j).text = "" if pd.isna(val) else f"{float(val):.0f}"
    prs.save(buf)


@app.callback(
    Output("dl-ppt-out", "data"),
    Input("dl-ppt", "n_clicks"),
    State("responses-store", "data"),
    prevent_initial_call=True,
)
def download_ppt(_, data):
    """
    Download the user's responses as a PPTX file.

    Args:
        _ (int): Click count of the "Download PPTX" button.
        data (dict): The user's responses, as stored in the "responses-store".

    Returns:
        dcc.SendData: A dcc.SendData object containing the PPTX data.
    """
    if not data:
        raise dash.exceptions.PreventUpdate
    return dcc.send_bytes(
        lambda b: _write_ppt_bytes(b, data), "Data_Maturity_Assessment.pptx"
    )


def _img_from_fig(fig, width=720, height=420, scale=2):
    # Requires kaleido installed
    """
    Convert a plotly figure to a PNG image bytes buffer.

    Args:
        fig (plotly.graph_objects.Figure): The figure to convert.
        width (int, optional): Image width in pixels. Defaults to 720.
        height (int, optional): Image height in pixels. Defaults to 420.
        scale (int, optional): Image resolution multiplier. Defaults to 2.

    Returns:
        io.BytesIO: A bytes buffer containing the PNG image data.
    """
    png_bytes = pio.to_image(fig, format="png", width=width, height=height, scale=scale)
    return io.BytesIO(png_bytes)


def _write_pdf_bytes(buf, data, theme):
    doc = SimpleDocTemplate(
        buf, pagesize=A4, leftMargin=16, rightMargin=16, topMargin=16, bottomMargin=16
    )
    styles = getSampleStyleSheet()
    story = []

    story += [
        Paragraph("<b>Data Maturity & Governance Assessment</b>", styles["Title"]),
        Spacer(1, 8),
    ]
    story += [
        Paragraph(
            f"Organization: {data.get('org','')}&nbsp;&nbsp;&nbsp; "
            f"Assessor: {data.get('assessor','')}&nbsp;&nbsp;&nbsp; "
            f"Sector: {data.get('sector','')}",
            styles["Normal"],
        ),
        Spacer(1, 10),
    ]
    story += [
        Paragraph(
            f"<b>Overall Maturity:</b> {data.get('overall',0):.1f}%", styles["Heading3"]
        ),
        Spacer(1, 8),
    ]

    # Domain Scores table with fixed widths
    d_scores = data.get("domain_scores", {})
    tbl_data = [["Domain", "Score (%)"]] + [
        [dom, f"{float(d_scores.get(dom, 0.0)):.1f}"] for dom in DOMAINS
    ]
    avail = A4[0] - 72
    col0 = 200
    col1 = avail - col0
    tbl = Table(tbl_data, colWidths=[col0, col1], hAlign="LEFT")
    tbl.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e9ebf3")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0b1020")),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                ("TOPPADDING", (0, 0), (-1, 0), 6),
            ]
        )
    )
    story += [
        Paragraph("<b>Domain Scores</b>", styles["Heading3"]),
        Spacer(1, 6),
        tbl,
        Spacer(1, 12),
    ]

    # Charts (PNG via kaleido)
    figs = [
        ("Maturity Radar", radar_figure(d_scores, theme)),
        ("Domain Bar", bar_figure(d_scores, theme)),
        (
            "Control Coverage Heatmap",
            heatmap_figure(pd.DataFrame(data.get("coverage_matrix", [])), theme),
        ),
    ]
    for title, fig in figs:
        story += [Paragraph(f"<b>{title}</b>", styles["Heading3"]), Spacer(1, 6)]
        img_buf = _img_from_fig(fig, width=520, height=320, scale=2)
        story += [RLImage(img_buf, width=520, height=320), Spacer(1, 12)]

    # Recommendations
    recs = data.get("recommendations", [])
    if recs:
        bullets = ListFlowable(
            [
                ListItem(
                    Paragraph(
                        f"[{r['domain']}] {r['action']} (current: {r['score']:.0f}%)",
                        styles["Normal"],
                    )
                )
                for r in recs
            ],
            bulletType="bullet",
        )
        story += [
            Paragraph("<b>Top Recommended Actions</b>", styles["Heading3"]),
            Spacer(1, 6),
            bullets,
            Spacer(1, 12),
        ]

    # Coverage table (Framework × Domain)
    mat = pd.DataFrame(data.get("coverage_matrix", []))
    if not mat.empty:
        pv = mat.pivot_table(
            index="Framework", columns="Domain", values="Coverage", aggfunc="mean"
        ).reindex(index=FRAMEWORK_ROWS, columns=DOMAINS)
        cov_data = [["Framework"] + list(pv.columns)]
        for fw in pv.index:
            row = [fw] + [
                ("" if pd.isna(pv.loc[fw, dom]) else f"{float(pv.loc[fw, dom]):.0f}%")
                for dom in pv.columns
            ]
            cov_data.append(row)
        fw_col = 200
        dcol = (avail - fw_col) / max(1, len(DOMAINS))
        cov_tbl = Table(
            cov_data, colWidths=[fw_col] + [dcol] * len(DOMAINS), hAlign="LEFT"
        )
        cov_tbl.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e9ebf3")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0b1020")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("ALIGN", (1, 1), (-1, -1), "CENTER"),
                    ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                    ("TOPPADDING", (0, 0), (-1, 0), 6),
                ]
            )
        )
        story += [
            Paragraph("<b>Control Coverage (Table)</b>", styles["Heading3"]),
            Spacer(1, 6),
            cov_tbl,
        ]

    doc.build(story)


@app.callback(
    Output("dl-pdf-out", "data"),
    Input("dl-pdf", "n_clicks"),
    State("responses-store", "data"),
    State("theme-store", "data"),
    prevent_initial_call=True,
)
def download_pdf(_, data, theme):
    """
    Download the user's responses as a PDF file.

    Args:
        _ (int): Click count of the "Download PDF" button.
        data (dict): The user's responses, as stored in the "responses-store".
        theme (str, optional): The theme to apply to the PDF (light or dark). Defaults to "light".

    Returns:
        dcc.SendData: A dcc.SendData object containing the PDF data.
    """
    if not data:
        raise dash.exceptions.PreventUpdate
    return dcc.send_bytes(
        lambda b: _write_pdf_bytes(b, data, theme or "light"),
        "Data_Maturity_Assessment.pdf",
    )


# UX: switch to results after computing
@app.callback(
    Output("tabs", "value"),
    Input("submit-assessment", "n_clicks"),
    prevent_initial_call=True,
)
def switch_to_results(n):
    """
    Switch the app to the "Results" tab after submitting the assessment.

    Args:
        n (int): The number of times the "Submit Assessment" button has been clicked.

    Returns:
        str: The ID of the tab to switch to.

    Raises:
        dash.exceptions.PreventUpdate: If the input is not a number.
    """
    if n:
        return "tab-results"
    raise dash.exceptions.PreventUpdate


# Theme toggle -> update page class and store
@app.callback(
    Output("page-root", "className"),
    Output("theme-store", "data"),
    Input("theme-switch", "on"),
)
def apply_theme(is_on):
    """
    Toggle the page theme class and store the current theme value.

    Args:
        is_on (bool): The on/off state of the theme switch.

    Returns:
        tuple: A pair of (page class name, theme name).
    """
    theme = "dark" if is_on else "light"
    return f"page theme-{theme}", theme


# ---------- Main -------------------
if __name__ == "__main__":
    app.run(debug=False)
